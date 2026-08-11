"""文書内画像の抽出とS3保存のテスト。

見ているのは3層:
  - app.parsers  … ファイルのバイト列 → 画像（PDFはページ画像、xlsx/pptxは貼り絵）
  - app.storage  … 画像のS3キーの形
  - app.ingest   … 抽出の振り分けと、画像チャンクとしての登録

サンプルファイルはテスト内でメモリ上に組み立て、実ファイルは置かない
（test_parsers.py と同じ方針）。解析ライブラリが無い環境ではその項目を skip する。
"""
from __future__ import annotations

import io
from unittest.mock import MagicMock, patch

import pytest

from app import parsers


def _png_bytes(width: int, height: int, color=(200, 30, 30)) -> bytes:
    """テスト用の単色PNGを作る。"""
    Image = pytest.importorskip("PIL.Image")
    buf = io.BytesIO()
    Image.new("RGB", (width, height), color).save(buf, format="PNG")
    return buf.getvalue()


# ---------------------------------------------------------------------------
# _from_bytes: 埋め込み画像のバイト列 → ExtractedImage
# ---------------------------------------------------------------------------


def test_from_bytes_keeps_original_png_bytes():
    """Claudeが扱える形式は再エンコードせず、原本のバイト列をそのまま残す。"""
    raw = _png_bytes(300, 200)

    image = parsers._from_bytes(raw, "図1")

    assert image is not None
    assert image.data == raw  # ★再エンコードしない★
    assert (image.ext, image.content_type) == (".png", "image/png")
    assert (image.width, image.height) == (300, 200)
    assert image.label == "図1"


def test_from_bytes_converts_unsupported_format_to_png():
    """BMP など Claude が受け取れない形式は PNG に変換して揃える。"""
    Image = pytest.importorskip("PIL.Image")
    buf = io.BytesIO()
    Image.new("RGB", (300, 200), (10, 10, 10)).save(buf, format="BMP")

    image = parsers._from_bytes(buf.getvalue(), "図1")

    assert image is not None
    assert (image.ext, image.content_type) == (".png", "image/png")
    with Image.open(io.BytesIO(image.data)) as converted:
        assert converted.format == "PNG"
        assert converted.size == (300, 200)


def test_from_bytes_drops_images_below_min_pixels(monkeypatch):
    """ロゴ・アイコン相当の小さい絵は捨てる（画像チャンクが埋め尽くされるのを防ぐ）。"""
    monkeypatch.setattr(parsers, "IMAGE_MIN_PIXELS", 100)

    # 幅だけ足りないケースも捨てる（どちらかが閾値未満なら不採用）
    assert parsers._from_bytes(_png_bytes(50, 400), "ロゴ") is None
    assert parsers._from_bytes(_png_bytes(400, 50), "罫線") is None
    assert parsers._from_bytes(_png_bytes(100, 100), "図") is not None  # 閾値ちょうどは通す


def test_from_bytes_returns_none_for_unreadable_bytes():
    """画像として読めないバイト列は例外にせず None（1枚落ちても取り込みは続ける）。"""
    assert parsers._from_bytes(b"not an image at all", "壊れた図") is None


# ---------------------------------------------------------------------------
# PDF: ページ全体をレンダリングする
# ---------------------------------------------------------------------------


def test_extract_pdf_images_renders_one_image_per_page():
    """埋め込み画像の有無に関係なく、1ページ = 1画像で返る。"""
    pytest.importorskip("pypdfium2")
    Image = pytest.importorskip("PIL.Image")

    # Pillow で2ページのPDFを作る（ページ内容はベクタでもラスタでも、
    # ページ画像化なので必ず取れる、という前提を確かめるためのもの）
    buf = io.BytesIO()
    page1 = Image.new("RGB", (600, 400), (255, 255, 255))
    page2 = Image.new("RGB", (600, 400), (0, 0, 255))
    page1.save(buf, format="PDF", save_all=True, append_images=[page2])

    images = parsers.extract_pdf_images(buf.getvalue())

    assert len(images) == 2
    assert [i.label for i in images] == ["1ページ目", "2ページ目"]
    assert all(i.ext == ".png" and i.content_type == "image/png" for i in images)
    with Image.open(io.BytesIO(images[0].data)) as rendered:
        assert rendered.format == "PNG"
    # スケール(既定2.0)を掛けた分だけ元のページより大きく描かれる
    assert images[0].width > 600


def test_extract_pdf_images_respects_max_per_doc(monkeypatch):
    """ページ数の多いPDFでも上限で打ち切る（S3とDBが膨らむのを止める）。"""
    pytest.importorskip("pypdfium2")
    Image = pytest.importorskip("PIL.Image")
    monkeypatch.setattr(parsers, "IMAGE_MAX_PER_DOC", 2)

    buf = io.BytesIO()
    first = Image.new("RGB", (200, 200), (255, 255, 255))
    rest = [Image.new("RGB", (200, 200), (255, 255, 255)) for _ in range(4)]
    first.save(buf, format="PDF", save_all=True, append_images=rest)

    assert len(parsers.extract_pdf_images(buf.getvalue())) == 2


def test_extract_pdf_images_returns_empty_on_broken_pdf():
    """壊れたPDFは例外を投げず空リスト（テキストだけで取り込みを続ける）。"""
    pytest.importorskip("pypdfium2")
    assert parsers.extract_pdf_images(b"this is definitely not a pdf") == []


def test_render_page_closes_the_bitmap_after_materializing_the_png():
    """★描画バッファを解放し、かつ解放の順番を守る★（レビュー指摘）

    PDFium のビットマップはC側の確保で、ページ1枚でも数MB。GC任せにすると
    大きなPDFでピークが跳ねる。一方 to_pil() はビットマップとメモリを共有する
    ので、先にビットマップを閉じると解放済み領域を読むことになる。
    PNG化 → PIL → ビットマップ の順であることを固定する。
    """
    closed: list[str] = []

    class FakePil:
        width, height = 100, 200

        def save(self, buf, format):
            assert "bitmap" not in closed, "PNG化より先にビットマップを閉じている"
            buf.write(b"png-bytes")

        def close(self):
            closed.append("pil")

    class FakeBitmap:
        def to_pil(self):
            return FakePil()

        def close(self):
            assert "pil" in closed, "PIL より先にビットマップを閉じている"
            closed.append("bitmap")

    class FakePage:
        def render(self, scale):
            return FakeBitmap()

    image = parsers._render_page({0: FakePage()}, 0)

    assert closed == ["pil", "bitmap"]
    assert image is not None
    assert image.data == b"png-bytes"
    assert (image.width, image.height) == (100, 200)


def test_render_page_closes_the_bitmap_even_when_rendering_fails():
    """途中で失敗しても確保済みのバッファは解放する。"""
    closed: list[str] = []

    class FakeBitmap:
        def to_pil(self):
            raise RuntimeError("描画に失敗")

        def close(self):
            closed.append("bitmap")

    class FakePage:
        def render(self, scale):
            return FakeBitmap()

    assert parsers._render_page({0: FakePage()}, 0) is None
    assert closed == ["bitmap"]


# ---------------------------------------------------------------------------
# XLSX / PPTX: 貼られた画像を取り出す
# ---------------------------------------------------------------------------


def test_extract_xlsx_images_picks_up_embedded_picture():
    openpyxl = pytest.importorskip("openpyxl")
    pytest.importorskip("PIL.Image")
    from openpyxl.drawing.image import Image as XlsxImage

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "売上"
    ws.add_image(XlsxImage(io.BytesIO(_png_bytes(400, 300))), "B2")
    buf = io.BytesIO()
    wb.save(buf)

    images = parsers.extract_xlsx_images(buf.getvalue())

    assert len(images) == 1
    assert "売上" in images[0].label  # どのシート由来かが分かる名前になっている
    assert (images[0].width, images[0].height) == (400, 300)


def test_extract_xlsx_images_returns_empty_when_no_picture():
    openpyxl = pytest.importorskip("openpyxl")

    wb = openpyxl.Workbook()
    wb.active["A1"] = "画像なし"
    buf = io.BytesIO()
    wb.save(buf)

    assert parsers.extract_xlsx_images(buf.getvalue()) == []


def test_extract_pptx_images_picks_up_picture_with_slide_label():
    pptx = pytest.importorskip("pptx")
    pytest.importorskip("PIL.Image")
    from pptx.util import Inches

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 6 = 空白レイアウト
    slide.shapes.add_picture(
        io.BytesIO(_png_bytes(500, 400)), Inches(1), Inches(1)
    )
    buf = io.BytesIO()
    prs.save(buf)

    images = parsers.extract_pptx_images(buf.getvalue())

    assert len(images) == 1
    assert images[0].label.startswith("Slide 1")
    assert (images[0].width, images[0].height) == (500, 400)


def test_extract_pptx_images_ignores_textboxes():
    """テキストボックスだけのスライドからは画像が出ない（図形の取り違え防止）。"""
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "画像のないスライド"
    buf = io.BytesIO()
    prs.save(buf)

    assert parsers.extract_pptx_images(buf.getvalue()) == []


def test_image_extractors_registry_maps_expected_extensions():
    assert set(parsers.IMAGE_EXTRACTORS) == {".pdf", ".xlsx", ".pptx"}
    assert parsers.IMAGE_EXTRACTORS[".pdf"] is parsers.extract_pdf_images


# ---------------------------------------------------------------------------
# storage.image_key: S3キーの形
# ---------------------------------------------------------------------------


def test_image_key_is_namespaced_and_zero_padded():
    storage = pytest.importorskip("app.storage")

    # 原本ファイル（キー=出典名）と名前空間を分ける + 辞書順=ページ順にする
    assert storage.image_key("決算.pdf", 1, ".png") == "images/決算.pdf/0001.png"
    assert storage.image_key("決算.pdf", 10, ".png") > storage.image_key("決算.pdf", 2, ".png")


# ---------------------------------------------------------------------------
# ingest.extract_images: 拡張子ごとの振り分け
# ---------------------------------------------------------------------------


def test_extract_images_returns_empty_for_text_file():
    ingest = pytest.importorskip("app.ingest")
    assert ingest.extract_images("memo.txt", b"just text") == []


def test_extract_images_delegates_by_extension():
    ingest = pytest.importorskip("app.ingest")
    marker = [parsers.ExtractedImage(b"x", ".png", "image/png", "1ページ目", 10, 10)]

    with patch.dict(parsers.IMAGE_EXTRACTORS, {".pdf": lambda data: marker}):
        assert ingest.extract_images("doc.PDF", b"%PDF-dummy") == marker  # 大文字も拾う


def test_extract_images_disabled_by_flag(monkeypatch):
    """EXTRACT_IMAGES=false なら拡張子によらず抽出しない（機能ごと止めるスイッチ）。"""
    ingest = pytest.importorskip("app.ingest")
    monkeypatch.setattr(ingest, "EXTRACT_IMAGES", False)

    with patch.dict(parsers.IMAGE_EXTRACTORS, {".pdf": lambda data: 1 / 0}):
        assert ingest.extract_images("doc.pdf", b"%PDF-dummy") == []


# ---------------------------------------------------------------------------
# ingest.store_images: S3保存 → 画像チャンク登録
# ---------------------------------------------------------------------------


def _image(label: str) -> "parsers.ExtractedImage":
    return parsers.ExtractedImage(b"bytes", ".png", "image/png", label, 400, 300)


def _mock_conn(next_index: int = 5):
    """get_conn() のモック。(get_conn, conn, cursor) を返す。"""
    conn = MagicMock()
    conn.execute.return_value.fetchone.return_value = [next_index]
    get_conn = MagicMock()
    get_conn.return_value.__enter__.return_value = conn
    return get_conn, conn, conn.cursor.return_value.__enter__.return_value


# INSERT の列並び（store_images）。位置で読むのでここに寄せておく。
_COL = {
    "document_id": 0,
    "chunk_index": 1,
    "content": 2,
    "context": 3,
    "content_nouns": 4,
    "embedding": 5,
    "image_embedding": 6,
    "image_path": 7,
}


def _col(rows: list[tuple], name: str) -> list:
    return [r[_COL[name]] for r in rows]


def test_store_images_saves_to_s3_and_inserts_image_chunks(monkeypatch):
    ingest = pytest.importorskip("app.ingest")
    get_conn, conn, cur = _mock_conn(next_index=5)
    saved: list[tuple[str, bytes, str]] = []

    monkeypatch.setattr(ingest.storage, "is_enabled", lambda: True)
    monkeypatch.setattr(
        ingest.storage,
        "save_bytes",
        lambda key, data, ct: (saved.append((key, data, ct)), True)[1],
    )
    monkeypatch.setattr(ingest, "get_conn", get_conn)

    count = ingest.store_images(
        7, "決算.pdf", [_image("1ページ目"), _image("2ページ目")], index_method="none"
    )

    assert count == 2
    # S3キーは1始まりの連番（storage.image_key）
    assert [key for key, _, _ in saved] == [
        "images/決算.pdf/0001.png",
        "images/決算.pdf/0002.png",
    ]
    assert saved[0][2] == "image/png"

    # 画像チャンクは chunk_index をテキストチャンクの続き番号にする
    rows = cur.executemany.call_args.args[1]
    assert _col(rows, "chunk_index") == [5, 6]
    assert _col(rows, "image_path") == [
        "images/決算.pdf/0001.png",
        "images/決算.pdf/0002.png",
    ]
    assert "1ページ目" in rows[0][_COL["content"]]  # content は NOT NULL なのでラベル
    assert _col(rows, "context") == ["1ページ目", "2ページ目"]  # 由来は必ず残す

    # 再取り込みで同じ絵が二重に積まれないよう、先に既存の画像チャンクを消す
    deletes = [
        c.args[0] for c in conn.execute.call_args_list if "DELETE" in str(c.args[0])
    ]
    assert deletes and "image_path IS NOT NULL" in deletes[0]


def test_store_images_skips_rows_whose_s3_save_failed(monkeypatch):
    """★S3に入らなかった画像はDBにも入れない★（実体の無いキーを指す行を残さない）。"""
    ingest = pytest.importorskip("app.ingest")
    get_conn, _, cur = _mock_conn(next_index=0)

    monkeypatch.setattr(ingest.storage, "is_enabled", lambda: True)
    # 1枚目だけ保存成功、2枚目は失敗（S3障害・権限不足など）
    results = iter([True, False])
    monkeypatch.setattr(
        ingest.storage, "save_bytes", lambda key, data, ct: next(results)
    )
    monkeypatch.setattr(ingest, "get_conn", get_conn)

    count = ingest.store_images(
        7, "決算.pdf", [_image("1ページ目"), _image("2ページ目")], index_method="none"
    )

    assert count == 1
    rows = cur.executemany.call_args.args[1]
    assert _col(rows, "image_path") == ["images/決算.pdf/0001.png"]


def test_store_images_indexes_only_images_that_reached_s3(monkeypatch):
    """★索引作成はS3保存の後★。保存できなかった画像にAPIコストを払わない。"""
    ingest = pytest.importorskip("app.ingest")
    get_conn, _, _ = _mock_conn(next_index=0)
    indexed: list[list] = []

    monkeypatch.setattr(ingest.storage, "is_enabled", lambda: True)
    results = iter([False, True])  # 1枚目は保存失敗
    monkeypatch.setattr(
        ingest.storage, "save_bytes", lambda key, data, ct: next(results)
    )
    monkeypatch.setattr(ingest, "get_conn", get_conn)
    monkeypatch.setattr(
        ingest,
        "build_image_index",
        lambda source, images, method: (
            indexed.append([i.label for i in images]),
            [
                {
                    "content": "x",
                    "context": i.label,
                    "content_nouns": None,
                    "embedding": None,
                    "image_embedding": None,
                }
                for i in images
            ],
        )[1],
    )

    ingest.store_images(7, "決算.pdf", [_image("1ページ目"), _image("2ページ目")])

    assert indexed == [["2ページ目"]]  # 保存できた1枚だけが索引作成にかかる


def test_store_images_does_nothing_without_s3(monkeypatch):
    """S3未設定なら画像は扱わない（DBにも触らない）。"""
    ingest = pytest.importorskip("app.ingest")
    get_conn, _, _ = _mock_conn()

    monkeypatch.setattr(ingest.storage, "is_enabled", lambda: False)
    monkeypatch.setattr(ingest, "get_conn", get_conn)

    assert ingest.store_images(7, "決算.pdf", [_image("1ページ目")]) == 0
    get_conn.assert_not_called()


def test_store_images_with_no_images_does_not_touch_db(monkeypatch):
    ingest = pytest.importorskip("app.ingest")
    get_conn, _, _ = _mock_conn()
    monkeypatch.setattr(ingest, "get_conn", get_conn)

    assert ingest.store_images(7, "決算.pdf", []) == 0
    get_conn.assert_not_called()


def test_store_images_skips_db_when_every_save_failed(monkeypatch):
    """全滅したときは DELETE も走らせない（既存の画像チャンクを巻き添えにしない）。"""
    ingest = pytest.importorskip("app.ingest")
    get_conn, _, _ = _mock_conn()

    monkeypatch.setattr(ingest.storage, "is_enabled", lambda: True)
    monkeypatch.setattr(ingest.storage, "save_bytes", lambda key, data, ct: False)
    monkeypatch.setattr(ingest, "get_conn", get_conn)

    assert ingest.store_images(7, "決算.pdf", [_image("1ページ目")]) == 0
    get_conn.assert_not_called()
