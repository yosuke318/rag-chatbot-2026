"""app.parsers（ファイル→テキスト抽出）のテスト。

解析ライブラリ（openpyxl / python-pptx / pypdf）が入っていない環境では
その項目を skip する（本番イメージには入っているのでそこでは実行される）。
サンプルファイルはテスト内でメモリ上に組み立て、実ファイルは置かない。
"""
import io

import pytest

from app.parsers import PARSERS, parse_pdf, parse_pptx, parse_xlsx


def test_parsers_registry_maps_expected_extensions():
    assert set(PARSERS) == {".pdf", ".xlsx", ".pptx"}
    assert PARSERS[".xlsx"] is parse_xlsx


def test_parse_xlsx_extracts_sheet_title_and_cells():
    openpyxl = pytest.importorskip("openpyxl")

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "勤怠"
    ws["A1"], ws["B1"] = "名前", "残業時間"
    ws["A2"], ws["B2"] = "太郎", 30
    buf = io.BytesIO()
    wb.save(buf)

    text = parse_xlsx(buf.getvalue())

    assert "# 勤怠" in text           # シート名は文脈として残す
    assert "名前\t残業時間" in text    # 行はタブ区切り
    assert "太郎\t30" in text          # 数値は文字列化される（None セルは落ちる）


def test_parse_pptx_extracts_slide_marker_and_text():
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 6 = 空白レイアウト
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "就業規則の概要"
    buf = io.BytesIO()
    prs.save(buf)

    text = parse_pptx(buf.getvalue())

    assert "# Slide 1" in text
    assert "就業規則の概要" in text


def test_parse_pdf_raises_valueerror_on_garbage():
    pytest.importorskip("pypdf")
    # PDFとして壊れたバイト列は ValueError にまとめられる（呼び出し側で400にする契約）
    with pytest.raises(ValueError):
        parse_pdf(b"this is definitely not a pdf")
