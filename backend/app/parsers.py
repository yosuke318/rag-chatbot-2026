"""ファイル → プレーンテキスト / 画像 の抽出（PDF / XLSX / PPTX）。

ingest.extract_text / ingest.extract_images から拡張子ごとに呼ばれる。
テキスト側は体裁を捨てて本文だけを取り出すことに徹し、画像側は
「原本として残す価値のある絵」を取り出す（図表・チャート・写真）。

方針:
  - 重いライブラリ（pypdf / pypdfium2 / openpyxl / python-pptx / Pillow）は
    関数内で遅延importする。テキスト貼り付けしか使わない構成では読み込まれない
    ようにするため（storage の boto3 と同じ方針）。
  - 解析に失敗したら ValueError を投げる。呼び出し側(/ingest-file)は400にまとめる。
    ただし★画像抽出は失敗しても例外を投げない★（後述の extract_*_images）。
  - スキャンPDFのように文字が取れないファイルは空文字を返す。呼び出し側が
    「本文が空」として弾く。
"""
from __future__ import annotations

import io
import logging
from dataclasses import dataclass

from app.config import IMAGE_MAX_PER_DOC, IMAGE_MIN_PIXELS, PDF_RENDER_SCALE

logger = logging.getLogger(__name__)


def parse_pdf(data: bytes) -> str:
    """PDFの各ページからテキストを抽出して連結する。

    テキストレイヤを持たない（スキャン画像の）PDFは空になる。OCRは範囲外。
    パスワード無しで開ける暗号化PDFは開く（空パスワードで復号を試みる）。
    """
    from pypdf import PdfReader
    from pypdf.errors import PdfReadError

    try:
        reader = PdfReader(io.BytesIO(data))
        if reader.is_encrypted:
            # 空パスワードで開けるものだけ対応（本当に保護されたものは失敗させる）
            reader.decrypt("")
        pages = [(page.extract_text() or "").strip() for page in reader.pages]
    except PdfReadError as exc:
        raise ValueError(f"PDFを解析できませんでした: {exc}") from exc
    # ページ境界は空行で区切る（チャンク分割時に文がまたがりにくくする）
    return "\n\n".join(p for p in pages if p)


def parse_xlsx(data: bytes) -> str:
    """XLSXの各シートのセル値を、シートごとにタブ区切りで並べる。

    data_only=True で数式は「計算済みの値」を取る（式そのものは検索に無意味なため）。
    read_only=True で大きなブックでもメモリを抑える。
    """
    from openpyxl import load_workbook

    try:
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    except Exception as exc:  # openpyxl は種々の例外を投げるためまとめて包む
        raise ValueError(f"XLSXを解析できませんでした: {exc}") from exc

    lines: list[str] = []
    try:
        for ws in wb.worksheets:
            lines.append(f"# {ws.title}")  # シート名は文脈になるので残す
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    lines.append("\t".join(cells))
    finally:
        wb.close()
    return "\n".join(lines)


def parse_pptx(data: bytes) -> str:
    """PPTXの各スライドから、図形内テキストと表のセルを抽出する。"""
    from pptx import Presentation

    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as exc:
        raise ValueError(f"PPTXを解析できませんでした: {exc}") from exc

    lines: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"# Slide {i}")
        for shape in slide.shapes:
            # テキストを持つ図形（タイトル・本文プレースホルダ・テキストボックス）
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    lines.append(text)
            # 表はセルを行ごとにタブ区切りで
            if shape.has_table:
                for r in shape.table.rows:
                    cells = [cell.text.strip() for cell in r.cells]
                    lines.append("\t".join(cells))
    return "\n".join(lines)


# 拡張子 → パーサ。ingest.extract_text がこの表を引く。
PARSERS = {
    ".pdf": parse_pdf,
    ".xlsx": parse_xlsx,
    ".pptx": parse_pptx,
}


# ---------------------------------------------------------------------------
# 画像の抽出（フェーズ5-1）
#
# テキスト抽出との一番の違いは「失敗しても取り込みを止めない」こと。本文が
# 取れなければ登録する意味が無いので ValueError で止めるが、図が1枚取れなくても
# 文書としては成立する。そのため下の extract_*_images は例外を投げず、
# 取れたものだけを返して残りは警告ログに落とす（best-effort）。
# ---------------------------------------------------------------------------

# Claude にそのまま渡せる画像フォーマット（5-3 で原本画像を回答生成に載せる）。
# ここに無い形式（BMP/TIFF 等）は Pillow で PNG に変換してから保存する。
SUPPORTED_IMAGE_FORMATS = {
    "PNG": (".png", "image/png"),
    "JPEG": (".jpg", "image/jpeg"),
    "GIF": (".gif", "image/gif"),
    "WEBP": (".webp", "image/webp"),
}


@dataclass(frozen=True)
class ExtractedImage:
    """文書から取り出した画像1枚。S3に保存する原本そのもの。

    label: 「1ページ目」「Slide 3」のような由来の名前。どこから出てきた絵かを
      利用者に示すのと、5-2 でキャプションを作るときの手がかりに使う。
    """

    data: bytes
    ext: str           # S3キーに付ける拡張子（".png" 等）
    content_type: str  # S3のContentType。ダウンロード時に正しく開くために要る
    label: str
    width: int
    height: int


def _too_small(width: int, height: int) -> bool:
    """ロゴ・罫線・アイコンとみなして捨てるサイズか（config.IMAGE_MIN_PIXELS）。"""
    return width < IMAGE_MIN_PIXELS or height < IMAGE_MIN_PIXELS


def _from_bytes(raw: bytes, label: str) -> ExtractedImage | None:
    """埋め込み画像のバイト列を ExtractedImage にする。使えなければ None。

    Claude が扱える形式なら★バイト列をそのまま残す★（再エンコードすると
    JPEGの写真が無駄に太るうえ、原本という位置づけからも遠ざかる）。
    それ以外の形式のときだけ Pillow で PNG に変換する。
    """
    from PIL import Image, UnidentifiedImageError

    try:
        with Image.open(io.BytesIO(raw)) as img:
            fmt, width, height = img.format, img.width, img.height
            if _too_small(width, height):
                return None
            known = SUPPORTED_IMAGE_FORMATS.get(fmt or "")
            if known is not None:
                ext, content_type = known
                return ExtractedImage(raw, ext, content_type, label, width, height)
            # 未対応形式（BMP/TIFF等）は PNG に変換して揃える
            buf = io.BytesIO()
            img.convert("RGB").save(buf, format="PNG")
    except (UnidentifiedImageError, OSError, ValueError):
        logger.warning("画像を読み取れませんでした（スキップ）: %s", label, exc_info=True)
        return None
    return ExtractedImage(buf.getvalue(), ".png", "image/png", label, width, height)


def extract_pdf_images(data: bytes) -> list[ExtractedImage]:
    """PDFの各ページを1枚の画像にレンダリングして返す（1ページ = 1画像）。

    ★埋め込み画像を拾うのではなくページごと描画する★のが肝。Excel/PowerPoint
    から出力したPDFの図表はベクタ描画で、埋め込みラスタ画像としては存在しない
    ことが多く、埋め込み画像だけを拾う実装では図表がまるごと抜け落ちる。
    ページ丸ごとなら、ベクタ図・チャート・スキャン画像のいずれも確実に残る。

    描画は pypdfium2（PDFium のバインディング。Apache/BSD系でライセンス上も
    商用配布に向く。PyMuPDF は AGPL なので採らない）。
    """
    import pypdfium2

    try:
        # パスワード無しで開ける暗号化PDFは開く（parse_pdf の decrypt("") と揃える）
        pdf = pypdfium2.PdfDocument(io.BytesIO(data), password="")
    except Exception:
        logger.warning("PDFを画像化できませんでした（テキストのみで継続）", exc_info=True)
        return []

    images: list[ExtractedImage] = []
    try:
        for i in range(min(len(pdf), IMAGE_MAX_PER_DOC)):
            page = _render_page(pdf, i)
            if page is not None:
                images.append(page)
    finally:
        pdf.close()
    return images


def _render_page(pdf, index: int) -> ExtractedImage | None:
    """PDFの1ページを PNG にする。失敗したら None（そのページだけ飛ばす）。

    ★描画バッファは使い終わったら必ず閉じる★
      PDFium のビットマップはC側に確保され、ページ1枚でも数MBになる
      （144dpiのA4で約8MB）。GC任せにすると大きなPDFでピーク使用量が跳ねる。

      解放の順番が重要で、to_pil() が返す PIL 画像は★ビットマップと同じメモリを
      共有する★（コピーではない）。先にビットマップを閉じると解放済みメモリを
      読むことになるため、PNGへ書き出して中身を確定させてから、PIL → ビットマップ
      の順に閉じる。
    """
    bitmap = None
    pil = None
    try:
        bitmap = pdf[index].render(scale=PDF_RENDER_SCALE)
        pil = bitmap.to_pil()
        buf = io.BytesIO()
        pil.save(buf, format="PNG")  # ここでバッファから独立したバイト列になる
        width, height = pil.width, pil.height
    except Exception:
        logger.warning("%dページ目を画像化できませんでした", index + 1, exc_info=True)
        return None
    finally:
        if pil is not None:
            pil.close()
        if bitmap is not None:
            bitmap.close()
    return ExtractedImage(
        buf.getvalue(), ".png", "image/png", f"{index + 1}ページ目", width, height
    )


def extract_xlsx_images(data: bytes) -> list[ExtractedImage]:
    """XLSXの各シートに貼られた画像（グラフ画像・写真・図）を取り出す。

    read_only=True では画像が読み込まれないため、ここだけ通常モードで開く
    （テキスト抽出の parse_xlsx とはワークブックを共有しない）。

    ※Excelの「グラフ」オブジェクト（chart）はベクタなので画像としては出てこない。
      画像として貼られた図とスクリーンショットが対象。
    """
    from openpyxl import load_workbook

    try:
        wb = load_workbook(io.BytesIO(data))
    except Exception:
        logger.warning("XLSXから画像を取り出せませんでした（テキストのみで継続）", exc_info=True)
        return []

    images: list[ExtractedImage] = []
    try:
        for ws in wb.worksheets:
            # _images は openpyxl の内部属性。公開APIが無いためこれを使う
            # （バージョン差で消えても取り込みは止めない、という前提で getattr）。
            for n, img in enumerate(getattr(ws, "_images", []), start=1):
                if len(images) >= IMAGE_MAX_PER_DOC:
                    return images
                raw = _image_ref_bytes(img)
                if raw is None:
                    continue
                extracted = _from_bytes(raw, f"シート「{ws.title}」の画像{n}")
                if extracted is not None:
                    images.append(extracted)
    finally:
        wb.close()
    return images


def _image_ref_bytes(img) -> bytes | None:
    """openpyxl の画像オブジェクトからバイト列を取り出す。

    `ref` の中身はバージョンや読み込み経路で変わる（PIL Image だったり
    BytesIO だったりする）ので、ここで形の違いを吸収する。
    """
    ref = getattr(img, "ref", None)
    if isinstance(ref, bytes):
        return ref
    if hasattr(ref, "read"):  # BytesIO 等
        try:
            ref.seek(0)
            return ref.read()
        except Exception:
            return None
    if hasattr(ref, "save"):  # PIL.Image
        try:
            buf = io.BytesIO()
            ref.save(buf, format=ref.format or "PNG")
            return buf.getvalue()
        except Exception:
            return None
    return None


def extract_pptx_images(data: bytes) -> list[ExtractedImage]:
    """PPTXの各スライドに貼られた画像を取り出す（グループ図形の中も辿る）。"""
    from pptx import Presentation

    try:
        prs = Presentation(io.BytesIO(data))
    except Exception:
        logger.warning("PPTXから画像を取り出せませんでした（テキストのみで継続）", exc_info=True)
        return []

    images: list[ExtractedImage] = []
    for i, slide in enumerate(prs.slides, start=1):
        for n, raw in enumerate(_pptx_picture_bytes(slide.shapes), start=1):
            if len(images) >= IMAGE_MAX_PER_DOC:
                return images
            extracted = _from_bytes(raw, f"Slide {i} の画像{n}")
            if extracted is not None:
                images.append(extracted)
    return images


def _pptx_picture_bytes(shapes) -> list[bytes]:
    """図形の並びから画像のバイト列を集める。グループ図形は中を再帰で辿る。"""
    out: list[bytes] = []
    for shape in shapes:
        # グループの中に図が入っているスライドは珍しくないので降りる
        if getattr(shape, "shapes", None) is not None:
            out.extend(_pptx_picture_bytes(shape.shapes))
            continue
        try:
            # 図以外の図形では属性そのものが無く(AttributeError)、
            # 外部リンク画像（実体を持たない図）は ValueError になる。
            # どちらも「この図形には画像が無い」として飛ばす。
            out.append(shape.image.blob)
        except AttributeError:
            continue
        except Exception:
            logger.warning("PPTXの画像を読み出せませんでした（スキップ）", exc_info=True)
    return out


# 拡張子 → 画像抽出。ingest.extract_images がこの表を引く。
# ここに無い拡張子（.txt 等）は「画像を持たない」＝空リスト。
IMAGE_EXTRACTORS = {
    ".pdf": extract_pdf_images,
    ".xlsx": extract_xlsx_images,
    ".pptx": extract_pptx_images,
}
